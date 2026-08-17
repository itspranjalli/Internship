"""Brand palette and layout primitives for the technical deep-dive deck.

The palette and the shape of these helpers are lifted from docs/demo/build_ppt.py
so the two decks look like one family. The difference: build_ppt.py keeps a
module-level ``prs`` singleton, which makes it un-importable. Here everything
hangs off a Deck object, so render_pptx.py can build a presentation without
side effects at import time.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(os.path.dirname(HERE), "img")
REPO = os.path.dirname(os.path.dirname(HERE))

# ---- palette (identical to build_ppt.py) ----
INK = RGBColor(0x0B, 0x0F, 0x17)
EDB = RGBColor(0x1F, 0x4E, 0x79)
EDB_DK = RGBColor(0x13, 0x35, 0x5E)
GREY = RGBColor(0x5A, 0x64, 0x72)
LIGHT = RGBColor(0xF1, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1A, 0x7F, 0x37)
AMBER = RGBColor(0xB5, 0x47, 0x08)
RED = RGBColor(0xB4, 0x23, 0x18)
# technical-deck additions
CODE_BG = RGBColor(0x11, 0x18, 0x27)
CODE_FG = RGBColor(0xD7, 0xE3, 0xF4)
CODE_ACC = RGBColor(0x7F, 0xB0, 0xE0)
RULE = RGBColor(0xD7, 0xDE, 0xE8)
MONO = "Consolas"
SANS = "Calibri"

FOOTER = ("EDB RIS(C) Grant Claim Automation  ·  ST Engineering HR  ·  "
          "S26-10249-RIS(C)  ·  Technical deep-dive")


class Deck:
    """A 16:9 presentation plus the layout helpers the renderers need."""

    def __init__(self, footer=FOOTER, prs=None):
        # ``prs`` lets a caller render into an existing deck (build_final_deck.py
        # fills the hand-authored template) instead of a fresh presentation.
        self.prs = prs or Presentation()
        if prs is None:
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        self.SW = self.prs.slide_width
        self.SH = self.prs.slide_height
        self.BLANK = self.prs.slide_layouts[6]
        self.footer_text = footer
        # title-only header: matches the deep-dive slides that were pasted into
        # the final template with their kicker and accent rule stripped
        self.plain_header = False

    @classmethod
    def open(cls, path, footer=None):
        return cls(footer=footer, prs=Presentation(path))

    # ---- primitives ----
    def slide(self, bg=WHITE):
        s = self.prs.slides.add_slide(self.BLANK)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
        return s

    def box(self, slide, l, t, w, h, fill=None, line=None,
            shape=MSO_SHAPE.RECTANGLE, line_w=Pt(1)):
        shp = slide.shapes.add_shape(shape, int(l), int(t), int(w), int(h))
        shp.shadow.inherit = False
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = line_w
        return shp

    def text(self, slide, l, t, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space=6, font=SANS):
        """runs: list of paragraphs; each paragraph = list of (text, size, bold, color)."""
        tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, para in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space)
            for (txt, size, bold, color) in para:
                r = p.add_run()
                r.text = txt
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.color.rgb = color
                r.font.name = font
        return tb

    def bullets(self, slide, items, l, t, w, h, size=16, color=INK, gap=10,
                marker="▸  ", marker_color=EDB):
        """items: str, or (lead, body) where lead renders bold."""
        tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            lead = it[0] if isinstance(it, tuple) else None
            body = it[1] if isinstance(it, tuple) else it
            r = p.add_run()
            r.text = marker
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = marker_color
            r.font.name = SANS
            if lead:
                r2 = p.add_run()
                r2.text = lead + "  "
                r2.font.size = Pt(size)
                r2.font.bold = True
                r2.font.color.rgb = INK
                r2.font.name = SANS
            r3 = p.add_run()
            r3.text = body
            r3.font.size = Pt(size)
            r3.font.color.rgb = color
            r3.font.name = SANS
        return tb

    def header(self, slide, kicker, title):
        if self.plain_header:
            self.text(slide, Inches(0.85), Inches(0.34), Inches(11.8), Inches(0.54),
                      [[(title, 26, True, INK)]], space=2)
            return
        self.box(slide, 0, 0, self.SW, Inches(1.35), fill=WHITE)
        self.box(slide, 0, Inches(1.35), self.SW, Pt(3), fill=EDB)
        self.box(slide, Inches(0.55), Inches(0.42), Inches(0.16), Inches(0.62), fill=EDB)
        self.text(slide, Inches(0.85), Inches(0.34), Inches(11.8), Inches(1.0), [
            [(kicker.upper(), 12, True, EDB)],
            [(title, 26, True, INK)],
        ], space=2)

    def footer(self, slide):
        if not self.footer_text:
            return
        self.text(slide, Inches(0.85), Inches(7.04), Inches(11.6), Inches(0.4),
                  [[(self.footer_text, 9, False, GREY)]])

    def img_fit(self, slide, path, l, t, max_w, max_h, border=True):
        if not os.path.isabs(path):
            path = os.path.join(IMG, path)
        if not os.path.exists(path):
            return None
        pic = slide.shapes.add_picture(path, int(l), int(t))
        ratio = min(max_w / pic.width, max_h / pic.height)
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
        pic.left = int(l + (max_w - pic.width) / 2)
        pic.top = int(t + (max_h - pic.height) / 2)
        if border:
            pic.line.color.rgb = RULE
            pic.line.width = Pt(1)
        return pic

    # ---- composite blocks used by the technical slides ----
    def code(self, slide, l, t, w, h, lines, size=12.5, title=None):
        """A dark monospace panel. lines: list of str, or (str, color)."""
        if title:
            self.text(slide, l, t - Inches(0.32), w, Inches(0.3),
                      [[(title, 11, True, GREY)]])
        self.box(slide, l, t, w, h, fill=CODE_BG)
        tb = slide.shapes.add_textbox(int(l + Inches(0.22)), int(t + Inches(0.14)),
                                      int(w - Inches(0.4)), int(h - Inches(0.28)))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, ln in enumerate(lines):
            txt, col = (ln if isinstance(ln, tuple) else (ln, CODE_FG))
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(2)
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = col
            r.font.name = MONO
        return tb

    def table(self, slide, l, t, w, cols, rows, col_w=None, size=11,
              head_fill=EDB, head_fg=WHITE, row_h=Inches(0.32),
              head_h=Inches(0.36), zebra=True, cell_colors=None):
        """A hand-drawn table. Native pptx tables fight the brand styling, and
        this keeps every cell a plain shape the presenter can restyle.
        cell_colors: optional {(row_idx, col_idx): RGBColor} for text colour.
        """
        n = len(cols)
        widths = col_w or [w / n] * n
        x = l
        for j, c in enumerate(cols):
            self.box(slide, x, t, widths[j], head_h, fill=head_fill)
            self.text(slide, x + Inches(0.08), t, widths[j] - Inches(0.1), head_h,
                      [[(c, size, True, head_fg)]], anchor=MSO_ANCHOR.MIDDLE)
            x += widths[j]
        y = t + head_h
        for i, row in enumerate(rows):
            if zebra and i % 2 == 1:
                self.box(slide, l, y, w, row_h, fill=LIGHT)
            x = l
            for j, cell in enumerate(row):
                col = (cell_colors or {}).get((i, j), INK)
                bold = j == 0
                self.text(slide, x + Inches(0.08), y, widths[j] - Inches(0.1), row_h,
                          [[(str(cell), size, bold, col)]], anchor=MSO_ANCHOR.MIDDLE)
                x += widths[j]
            self.box(slide, l, y + row_h, w, Pt(0.5), fill=RULE)
            y += row_h
        return y

    def stat(self, slide, l, t, w, h, value, label, note=None, color=EDB):
        self.box(slide, l, t, w, h, fill=LIGHT)
        self.box(slide, l, t, Pt(4), h, fill=color)
        self.text(slide, l + Inches(0.2), t + Inches(0.12), w - Inches(0.3), Inches(0.5),
                  [[(value, 24, True, color)]])
        self.text(slide, l + Inches(0.2), t + Inches(0.62), w - Inches(0.3), Inches(0.3),
                  [[(label, 11, True, INK)]])
        if note:
            self.text(slide, l + Inches(0.2), t + Inches(0.9), w - Inches(0.3), Inches(0.5),
                      [[(note, 9.5, False, GREY)]])

    def connector(self, slide, x1, y1, x2, y2, color=EDB, width=Pt(1.5),
                  arrow=True, dashed=False):
        """Straight connector. Elbows are drawn as two straight segments by the
        caller so the geometry stays reproducible in both renderers.
        """
        cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        int(x1), int(y1), int(x2), int(y2))
        cxn.line.color.rgb = color
        cxn.line.width = width
        if dashed:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            cxn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if arrow:
            _arrow_end(cxn)
        return cxn

    def save(self, path):
        self.prs.save(path)
        return path


def strip_placeholders(slide, empty_only=True):
    """Drop layout placeholders from a slide.

    python-pptx exposes no shape deletion, so the element is unlinked from its
    parent directly. ``empty_only`` keeps any placeholder that already carries
    text, so a hand-written slide is never silently blanked.
    """
    removed = 0
    for sh in list(slide.shapes):
        if not sh.is_placeholder:
            continue
        if empty_only and sh.has_text_frame and sh.text_frame.text.strip():
            continue
        sh._element.getparent().remove(sh._element)
        removed += 1
    return removed


def clear_slide(slide):
    """Remove every shape, placeholder or not — for replacing a pasted slide."""
    removed = 0
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)
        removed += 1
    return removed


def move_slide(prs, from_idx, to_idx):
    """Reposition a slide. python-pptx can only append, so reorder the id list."""
    ids = prs.slides._sldIdLst
    entries = list(ids)
    entry = entries[from_idx]
    ids.remove(entry)
    ids.insert(to_idx, entry)


def _arrow_end(cxn):
    """python-pptx exposes no arrowhead API; set it on the line element directly."""
    from pptx.oxml.ns import qn
    ln = cxn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = ln.makeelement(qn("a:tailEnd"), {})
        ln.append(tail)
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
