"""Generate the stakeholder presentation (docs/demo/EDB_System_Overview.pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")

INK = RGBColor(0x0B, 0x0F, 0x17)
EDB = RGBColor(0x1F, 0x4E, 0x79)
EDB_DK = RGBColor(0x13, 0x35, 0x5E)
GREY = RGBColor(0x5A, 0x64, 0x72)
LIGHT = RGBColor(0xF1, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1A, 0x7F, 0x37)
AMBER = RGBColor(0xB5, 0x47, 0x08)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    """runs: list of paragraphs; each paragraph = list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for (txt, size, bold, color) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def bullets(slide, items, l, t, w, h, size=16, color=INK, gap=10):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        lead = it[0] if isinstance(it, tuple) else None
        body = it[1] if isinstance(it, tuple) else it
        r = p.add_run(); r.text = "▸  "
        r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = EDB; r.font.name = "Calibri"
        if lead:
            r2 = p.add_run(); r2.text = lead + "  "
            r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = INK; r2.font.name = "Calibri"
        r3 = p.add_run(); r3.text = body
        r3.font.size = Pt(size); r3.font.color.rgb = color; r3.font.name = "Calibri"
    return tb


def header(slide, kicker, title):
    _box(slide, 0, 0, SW, Inches(1.35), fill=WHITE)
    _box(slide, 0, Inches(1.35), SW, Pt(3), fill=EDB)
    _box(slide, Inches(0.55), Inches(0.42), Inches(0.16), Inches(0.62), fill=EDB)
    _text(slide, Inches(0.85), Inches(0.34), Inches(11.8), Inches(1.0), [
        [(kicker.upper(), 12, True, EDB)],
        [(title, 28, True, INK)],
    ], space=2)


def img_fit(slide, path, l, t, max_w, max_h, border=True):
    if not os.path.exists(path):
        return None
    pic = slide.shapes.add_picture(path, l, t)
    ratio = min(max_w / pic.width, max_h / pic.height)
    pic.width = int(pic.width * ratio); pic.height = int(pic.height * ratio)
    pic.left = int(l + (max_w - pic.width) / 2)
    pic.top = int(t + (max_h - pic.height) / 2)
    if border:
        pic.line.color.rgb = RGBColor(0xD7, 0xDE, 0xE8); pic.line.width = Pt(1)
    return pic


def two_col_slide(kicker, title, lhead, litems, rhead, ritems, lcolor=AMBER, rcolor=GREEN):
    s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
    header(s, kicker, title)
    top = Inches(1.95); h = Inches(4.75); w = Inches(5.75)
    for x, head, items, col in ((Inches(0.85), lhead, litems, lcolor),
                                (Inches(6.95), rhead, ritems, rcolor)):
        _box(s, x, top, w, Inches(0.62), fill=col)
        _text(s, x, top, w, Inches(0.62), [[(head, 16, True, WHITE)]],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _box(s, x, top + Inches(0.62), w, h - Inches(0.62), fill=LIGHT)
        bullets(s, items, x + Inches(0.25), top + Inches(0.85), w - Inches(0.5),
                h - Inches(1.0), size=14, gap=11)
    _text(s, Inches(0.85), Inches(7.04), Inches(11.6), Inches(0.4),
          [[("EDB RIS(C) Grant Claim Automation  ·  ST Engineering HR  ·  Confidential", 9, False, GREY)]])
    return s


def flow_slide(kicker, title, steps, caption):
    """A horizontal chevron process ribbon with a sub-caption under each step."""
    from pptx.enum.shapes import MSO_SHAPE
    s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
    header(s, kicker, title)
    n = len(steps)
    w = Inches(2.45); overlap = Inches(0.35); step = w - overlap
    left = Inches(0.7); top = Inches(2.7); h = Inches(1.0)
    for i, (label, sub) in enumerate(steps):
        x = left + step * i
        ch = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x, top, w, h)
        ch.shadow.inherit = False
        ch.fill.solid(); ch.fill.fore_color.rgb = EDB if i % 2 == 0 else EDB_DK
        ch.line.fill.background()
        tf = ch.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"{i + 1}. {label}"
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
        # sub-caption under the chevron body
        cap = s.shapes.add_textbox(x + Inches(0.15), top + h + Inches(0.12), w - Inches(0.55), Inches(1.0))
        cap.text_frame.word_wrap = True
        cp = cap.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = sub
        cr.font.size = Pt(11); cr.font.color.rgb = GREY; cr.font.name = "Calibri"
    _box(s, Inches(0.9), Inches(5.6), Inches(11.5), Pt(2), fill=LIGHT)
    _text(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.9),
          [[(caption, 15, False, INK)]], align=PP_ALIGN.CENTER)
    _text(s, Inches(0.85), Inches(7.04), Inches(11.6), Inches(0.4),
          [[("EDB RIS(C) Grant Claim Automation  ·  ST Engineering HR  ·  Confidential", 9, False, GREY)]])
    return s


def content_slide(kicker, title, items, img=None, size=16):
    s = prs.slides.add_slide(BLANK); _bg(s, WHITE)
    header(s, kicker, title)
    if img:
        bullets(s, items, Inches(0.85), Inches(1.75), Inches(6.0), Inches(5.2), size=size)
        img_fit(s, os.path.join(IMG, img), Inches(7.1), Inches(1.75), Inches(5.7), Inches(5.2))
    else:
        bullets(s, items, Inches(0.95), Inches(1.95), Inches(11.4), Inches(5.0), size=size + 2, gap=14)
    _text(s, Inches(0.85), Inches(7.04), Inches(11.6), Inches(0.4),
          [[("EDB RIS(C) Grant Claim Automation  ·  ST Engineering HR  ·  Confidential", 9, False, GREY)]])
    return s


# ---- S1 Title ----
s = prs.slides.add_slide(BLANK); _bg(s, INK)
_box(s, 0, Inches(3.05), SW, Pt(3), fill=EDB)
_text(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.5),
      [[("EDB SUPPORT PACKAGE FOR AI COE  ·  RESEARCH INCENTIVE SCHEME (RIS(C))", 13, True, RGBColor(0x8F, 0xB2, 0xD8))]])
_text(s, Inches(0.9), Inches(1.6), Inches(11.6), Inches(1.6),
      [[("EDB Grant Claim Automation", 46, True, WHITE)],
       [("Manpower claims — prepared, verified, audit-ready", 22, False, RGBColor(0xC8, 0xD6, 0xE8))]], space=8)
_text(s, Inches(0.9), Inches(3.35), Inches(11.6), Inches(1.2), [
    [("ST Engineering  ·  Human Resources", 16, True, WHITE)],
    [("Application No. S26-10249-RIS(C)   |   Qualifying Period 01 Jan 2026 – 31 Dec 2028   |   17 participating entities", 13, False, RGBColor(0x9F, 0xB0, 0xC4))],
], space=6)
_text(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4),
      [[("Stakeholder briefing  ·  Confidential", 11, False, RGBColor(0x6E, 0x80, 0x96))]])

# ---- S2 Challenge ----
content_slide("The challenge", "Grant claims are high-stakes and manual", [
    ("Scale.", "17 participating entities, a three-year qualifying period, and a grant worth up to S$42 million — manpower only."),
    ("Manual & error-prone.", "Each claim is hand-built from timesheets, payroll registers and ECMF lists; one wrong rate or a missed exclusion can sink it."),
    ("Heavily audited.", "An external Public Accountant verifies under SSRS 4400 and samples at least 85% of the claimed value — every figure must trace to a source document."),
    ("Hard deadlines.", "Per-entity audits from Oct 2026, group submission to EDB by end Nov 2026, final audited claim within 183 days of the period end."),
])

# ---- S3 Solution ----
content_slide("The solution", "An HR-first portal that does the heavy lifting", [
    ("Upload, and it takes over.", "Drop the HR records; the system reads them, applies EDB's rules, and prepares the claim."),
    ("End to end.", "Eligibility checks → per-person calculation → the EDB submission pack — in one guided flow."),
    ("Private by design.", "Runs entirely on your machine; salary data never leaves the building."),
    ("Deterministic.", "The same inputs always produce identical outputs — a hard audit requirement."),
], img="shot_1_landing.png")

# ---- S3b Advantages ----
content_slide("Why this system", "The advantages at a glance", [
    ("Faster.", "Weeks of manual spreadsheet work per claim become minutes — upload to audit-ready pack."),
    ("Accurate.", "Rules applied the same way every time; reconciles to the cent and removes rate, cap and exclusion mistakes."),
    ("Audit-ready.", "Every figure is one click from its source cell — a smoother, faster, lower-cost SSRS 4400 audit."),
    ("Compliant.", "The 60% rate, 9-month cap, S$42m ceiling and reporting deadlines are enforced, not remembered."),
    ("Transparent.", "No black box: full workings are shown and nobody is silently dropped — every decision has a reason."),
    ("Secure & scalable.", "Runs 100% on-premise (salary data never leaves), and applies one consistent method across all 17 entities."),
], size=15)

# ---- S3c Before vs After ----
two_col_slide("The advantages", "Before vs after", "Manual today", [
    "Weeks of spreadsheet work per claim",
    "Easy to mis-apply the rate, cap or an exclusion",
    "Audit prep: hunting for source documents by hand",
    "Hard to keep 17 entities consistent",
    "A rejected claim risks delay or clawback",
    "Salary data scattered across files",
], "With this system", [
    "Minutes: upload → audit-ready submission pack",
    "Rules enforced automatically; cent-exact",
    "Every figure one click from its source cell",
    "One consistent method across all 17 entities",
    "Defensible, traceable, deterministic outputs",
    "All processing on-premise — nothing leaves",
])

# ---- S3d Workflow at a glance (process ribbon) ----
flow_slide("The process", "How a claim flows through the system", [
    ("Upload", "Drop the HR documents (timesheets, ECMF list, payroll, supporting evidence)."),
    ("Auto-check", "System tracks them against the required list and flags what's missing."),
    ("Eligibility", "Each person tested against EDB gates G1–G7."),
    ("Calculate", "Method A claim per person + internal cross-check."),
    ("Submission pack", "EDB template, SOE and issues list generated."),
    ("Audit", "Practitioner verifies under SSRS 4400; report goes to EDB."),
], "Upload once → the system checks, calculates and assembles the audit-ready pack. "
   "HR stays in control; every figure is traceable and nobody is silently dropped.")

# ---- S4 Workflow ----
content_slide("How it works", "Five guided, transparent steps", [
    ("Documents.", "Upload timesheets, the ECMF list and payroll; tick the supporting evidence."),
    ("Document check.", "Missing-document validation per company and per person."),
    ("Eligibility.", "Every person checked against EDB's criteria."),
    ("Claim amount.", "Each claim calculated and shown in full."),
    ("Submission pack.", "The audit-ready outputs, generated."),
    ("Verbose & re-runnable.", "You see exactly what it's doing; fix a document and re-run."),
], img="shot_2_analyzing.png")

# ---- S5 Trust ----
content_slide("Trust by design", "The AI assists — the maths is deterministic", [
    ("The model never computes a claim figure.", "All arithmetic is pure, auditable Python. The AI only reads documents and explains."),
    ("Full traceability.", "Every number links to its exact source cell — payslip, timesheet, ECMF list."),
    ("Nobody is silently dropped.", "Excluded and blocked people are always listed with the reason."),
    ("Proven.", "115 automated tests; reconciles to the cent against EDB's own worked example."),
])

# ---- S6 Eligibility ----
content_slide("Eligibility", "Seven EDB gates, applied to every person", [
    ("The checks.", "Citizen/PR · ECMF-validated · no other grant · salary ≥ S$5,000 · eligible R&D role · active in period · payslip present."),
    ("Clear verdicts.", "Qualifies · Not eligible · Needs a document — with the reason behind each."),
    ("Human in the loop.", "Borderline roles are flagged for HR to confirm, not auto-decided."),
], img="shot_4_eligibility.png")

# ---- S7 Calculation ----
content_slide("The calculation", "Every claim, shown in full", [
    ("EDB monthly pro-ration (Method A).", "Capped salary × month involved × time on project, summed × the support rate. This is the figure submitted."),
    ("Open any person.", "See the month-by-month working, each input linked to its source document."),
    ("Silent cross-check (Method B).", "An internal hours-based method runs in the background; differences are flagged, never hidden."),
], img="calc_1_simple.png")

# ---- S8 Scheme rules ----
content_slide("EDB scheme rules, encoded", "The support package, enforced automatically", [
    ("Support rate 60%", "of the qualifying monthly salary (EDB Support Package)."),
    ("Basic salary only.", "No CPF, bonus, AWS, allowances, COLA or airfare."),
    ("Floor S$5,000 · Cap S$20,000/month", "→ at most S$12,000 funded per person per month."),
    ("New hire vs upskill.", "New hires funded across the qualifying period; upskilling to PL3 funded for up to 9 months."),
    ("Ceiling S$42m, claim ≥ 3 months.", "Manpower-only grant cap, with a minimum claim period."),
], size=15)

# ---- S9 Grant & compliance ----
content_slide("Grant & compliance", "Ceiling, disbursement and deadlines — tracked", [
    ("Against the S$42m ceiling.", "Every submission is measured against the manpower grant cap."),
    ("70 / 30 disbursement gate.", "Up to 70% disbursed before completion; the final 30% on project completion + T&Cs."),
    ("Deadlines computed for you.", "12-month Practitioner's Report, 183-day final claim, progress updates, inspection rights."),
], img="grant_compliance.png")

# ---- S10 Evidence ----
content_slide("Intelligent document intake", "Upload anything reasonable — the system sorts it out", [
    ("Guided, in order.", "EDB output template and trainee list first, then timesheet / ECMF / payroll — a live tracker ticks each off and shows what's still missing."),
    ("Auto-categorised.", "Files are recognised by name; a leave report or CPF statement is never mistaken for the core timesheet."),
    ("Flexible formats.", "Synonym headers ('Emp No', 'Salary'), a Period/Pay-Date column, or one payslip per file with the month in its name — all parse and merge."),
    ("Clear reasons when rejected.", "If a file can't be read it explains why (missing sheet/column, wrong period) and how to fix it — never a stack trace."),
], size=15)

content_slide("Audit-ready evidence", "From any figure to its source, in one click", [
    ("One-click traceability.", "Click a number → the exact cell in the original payslip or timesheet, highlighted."),
    ("Three outputs, one per reader.", "EDB submission template · Statement of Expenditure (auditor) · Issues list (HR)."),
    ("Built for SSRS 4400.", "Designed for the auditor's ≥85% value sampling — every claimed dollar is traceable."),
], img="calc_evidence.png")

# ---- S10b SSRS 4400 explainer ----
content_slide("How claims are verified", "The audit explained — SSRS 4400", [
    ("Agreed-upon procedures — not an audit or a review.", "An ACRA-registered Public Accountant (the “Practitioner”) performs a set of procedures agreed with EDB and reports factual findings only — no opinion or assurance is given; EDB draws its own conclusions."),
    ("At least 85% of claimed value is sampled.", "The Practitioner tests ≥ 85% of the claimed value against the Statement of Expenditure (SOE), vouching each sampled figure to its source document."),
    ("Direct to EDB, on a cadence.", "The Practitioner’s Report goes straight from the Practitioner to EDB — at least every 12 months, with the final audited claim within 183 days of the qualifying-period end."),
    ("What we built for it.", "One-click traceability and the SOE output make this engagement fast and clean; outputs stay labelled “UnauditedClaim” until the Practitioner completes their procedures."),
], size=15)

# ---- S11 Assistant ----
content_slide("The built-in assistant", "Ask the claim anything, in plain English", [
    ("Natural-language Q&A.", "“Why is this person excluded?”, “How is the claim calculated?”, “What's the 9-month rule?”"),
    ("Recall across claims.", "“What documents do we have for E001?” fetches that person's saved document checklist from an earlier claim — with citations — even in a new session."),
    ("Local AI.", "Powered by a model running on ST Engineering's own DGX — private, no external calls."),
    ("Grounded.", "It explains and reasons, but never invents a figure — answers cite the source."),
], img="chat_model.png", size=14)

# ---- S12 Architecture ----
content_slide("Architecture & security", "Modern, on-premise, auditable", [
    ("Web app on a Python engine.", "A clean React interface over the deterministic calculation core."),
    ("Local AI + local data.", "Qwen model on the DGX; a per-employee SQLite store on the machine."),
    ("Nothing leaves the building.", "All salary and claim data processed on-premise; no internet calls."),
    ("Auditable AI.", "Every model call is temperature-0, schema-constrained, cached and logged."),
])

# ---- S13 Status ----
content_slide("Status & next steps", "POC complete — ready to pilot", [
    ("Built and verified.", "Full pipeline end-to-end on synthetic data; 115 automated tests passing."),
    ("Pending confirmation.", "Exact support rate & 9-month anchoring (Letter of Award); the auditor's final document list."),
    ("Next.", "Pilot on a real entity's data; connect to live payroll and timesheet exports; finalise outputs with the appointed auditor."),
])

# ---- S14 Closing ----
s = prs.slides.add_slide(BLANK); _bg(s, INK)
_box(s, 0, Inches(3.5), SW, Pt(3), fill=EDB)
_text(s, Inches(0.9), Inches(2.5), Inches(11.6), Inches(1.6),
      [[("Thank you", 44, True, WHITE)],
       [("Questions & discussion", 22, False, RGBColor(0xC8, 0xD6, 0xE8))]], space=10)
_text(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.4),
      [[("EDB Grant Claim Automation  ·  ST Engineering HR  ·  S26-10249-RIS(C)", 11, False, RGBColor(0x6E, 0x80, 0x96))]])

out = os.path.join(HERE, "EDB_System_Overview.pptx")
prs.save(out)
print("SAVED", out, "·", len(prs.slides._sldIdLst), "slides")
